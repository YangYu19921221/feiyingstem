#!/usr/bin/env python3
"""PPTX 构建库 — 品牌化幻灯片布局组件。

设计语言:活力橙 #FF6B35 / 阳光黄 #FFD23F / 天空蓝 #00D9FF / 草绿 #5FD35F / 米白 #FFF8F0。
16:9,中文标题黑体思源/微软雅黑,大字少字,给加盟商与学校看。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
from pathlib import Path

# 品牌色
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
YELLOW = RGBColor(0xFF, 0xD2, 0x3F)
BLUE   = RGBColor(0x00, 0xB4, 0xD8)   # 稍压暗的天空蓝,投影更清晰
SKY    = RGBColor(0x00, 0xD9, 0xFF)
GREEN  = RGBColor(0x35, 0xB8, 0x5F)
CREAM  = RGBColor(0xFF, 0xF8, 0xF0)
INK    = RGBColor(0x2A, 0x2A, 0x33)   # 主文字
GRAY   = RGBColor(0x6B, 0x6B, 0x78)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RED    = RGBColor(0xE5, 0x53, 0x53)
DARK   = RGBColor(0x1A, 0x1F, 0x2E)   # 深底封面

FONT = "微软雅黑"
IMG = Path(__file__).parent / "images"

SW, SH = Inches(13.333), Inches(7.5)


def new_deck():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _no_line(shape):
    shape.line.fill.background()


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def round_rect(slide, x, y, w, h, color, line=None, radius=0.12):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.5)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
            r._r.get_or_add_rPr().set(qn('w:eastAsia') if False else 'altLang', 'zh-CN')
    return tb


def _one(txt, size, color, bold=False):
    return [[(txt, size, color, bold)]]


def pic_cover(slide, path, x, y, w, h):
    """按 cover 方式裁剪填充图片到指定框(保持比例、裁边)。"""
    p = Path(path)
    if not p.exists():
        rect(slide, x, y, w, h, RGBColor(0xE8, 0xE8, 0xEE))
        return None
    iw, ih = Image.open(p).size
    tw, th = w, h
    box_ratio = tw / th
    img_ratio = iw / ih
    pic = slide.shapes.add_picture(str(p), x, y, w, h)
    if img_ratio > box_ratio:
        crop = (1 - box_ratio / img_ratio) / 2
        pic.crop_left = crop; pic.crop_right = crop
    else:
        crop = (1 - img_ratio / box_ratio) / 2
        pic.crop_top = crop; pic.crop_bottom = crop
    return pic


def pic_fit(slide, path, x, y, w, h, align_center=True):
    """完整放入(contain),不裁剪。"""
    p = Path(path)
    if not p.exists():
        return None
    iw, ih = Image.open(p).size
    img_ratio = iw / ih
    box_ratio = w / h
    if img_ratio > box_ratio:
        nw = w; nh = int(w / img_ratio)
    else:
        nh = h; nw = int(h * img_ratio)
    nx = x + (w - nw) // 2 if align_center else x
    ny = y + (h - nh) // 2 if align_center else y
    return slide.shapes.add_picture(str(p), nx, ny, nw, nh)


def pill(slide, x, y, w, h, txt, fill, txt_color, size=14, bold=True):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; _no_line(sp); sp.shadow.inherit = False
    try: sp.adjustments[0] = 0.5
    except Exception: pass
    tf = sp.text_frame; tf.word_wrap = False
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.color.rgb = txt_color; r.font.bold = bold; r.font.name = FONT
    return sp
