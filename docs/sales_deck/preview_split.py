#!/usr/bin/env python3
"""把整份 pptx 拆成每页单张 pptx,便于 qlmanage 逐页渲染预览。"""
import copy
import sys
from pathlib import Path
from pptx import Presentation

src = Path(sys.argv[1])
outdir = Path("/tmp/deckslides")
outdir.mkdir(exist_ok=True)
for f in outdir.glob("*.pptx"):
    f.unlink()

base = Presentation(str(src))
n = len(base.slides._sldIdLst)
for keep in range(n):
    prs = Presentation(str(src))
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i, sid in enumerate(ids):
        if i != keep:
            lst.remove(sid)
    prs.save(str(outdir / f"s{keep+1:02d}.pptx"))
print(f"split {n} slides -> {outdir}")
